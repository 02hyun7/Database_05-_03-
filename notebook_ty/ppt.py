
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

input_file = "erd_template.pptx"
output_file = "erd_template_visible.pptx"

prs = Presentation(input_file)

TITLE_FONT_SIZE = Pt(30)
ENTITY_FONT_SIZE = Pt(26)
LABEL_FONT_SIZE = Pt(20)
SHAPE_SCALE_FACTOR = 1.15
LINE_THICKNESS = Pt(2)

def scale_shape(shape, factor=1.15):
    try:
        cx = shape.left + shape.width / 2
        cy = shape.top + shape.height / 2
        new_w = int(shape.width * factor)
        new_h = int(shape.height * factor)
        shape.width = new_w
        shape.height = new_h
        shape.left = int(cx - new_w / 2)
        shape.top = int(cy - new_h / 2)
    except:
        pass

for slide in prs.slides:
    for shp in slide.shapes:
        if hasattr(shp, "text_frame") and shp.text_frame:
            tf = shp.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    txt = run.text.strip().lower()
                    if txt in {"er diagram", "ie notation", "example", "notations", "type"}:
                        run.font.size = TITLE_FONT_SIZE
                        run.font.bold = True
                    elif len(txt) <= 12:
                        run.font.size = ENTITY_FONT_SIZE
                        run.font.bold = True
                    else:
                        run.font.size = LABEL_FONT_SIZE
                        run.font.bold = True
                if len(para.text.strip()) <= 12:
                    para.alignment = PP_ALIGN.CENTER
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            scale_shape(shp, factor=SHAPE_SCALE_FACTOR)
        else:
            scale_shape(shp, factor=1.08)
            if hasattr(shp, "line"):
                try:
                    shp.line.width = LINE_THICKNESS
                except:
                    pass

prs.save(output_file)
print(f"Saved as {output_file}")
